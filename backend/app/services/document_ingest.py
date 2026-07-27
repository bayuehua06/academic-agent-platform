"""文档解析入库：粘贴 / 上传（md/docx/pdf/pptx/txt）→ raw_text。"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from pathlib import Path
from typing import Optional
from uuid import UUID, uuid4

from app.core.config import get_settings
from app.services.pandoc_service import pandoc_service

logger = logging.getLogger(__name__)
settings = get_settings()

ROLES = frozenset({"ASSESSMENT", "BACKGROUND", "OUTLINE", "SPECIFIC"})
SOURCE_TYPES = frozenset({"PASTE", "UPLOAD", "NOTEBOOKLM"})

PLACEHOLDER_MAX = 4000
# 单文档入库文本上限（超长 PDF 会议录等）
RAW_TEXT_STORE_MAX = 500_000


@dataclass
class IngestResult:
    """单次解析结果（尚未写库）。"""

    raw_text: str
    storage_path: Optional[str] = None
    original_filename: Optional[str] = None
    content_type: Optional[str] = None
    title: Optional[str] = None


def sanitize_extracted_text(text: str) -> str:
    """
    清理抽取文本，避免 PostgreSQL UTF-8 写入失败。

    PDF 抽取常混入 NUL(0x00) 与其它控制字符。
    """
    if not text:
        return ""
    # 去掉 NUL：asyncpg/PG 报 invalid byte sequence for encoding UTF8: 0x00
    cleaned = text.replace("\x00", "")
    # 保留换行/回车/制表符，去掉其余 C0 控制符
    return "".join(ch for ch in cleaned if ch in "\n\r\t" or ord(ch) >= 32)


def placeholder_summary(raw_text: str, max_len: int = PLACEHOLDER_MAX) -> str:
    """Phase 2 占位摘要；Phase 3 由 LLM 替换。"""
    text = (raw_text or "").strip()
    if not text:
        return ""
    if len(text) <= max_len:
        return text
    return (
        text[:max_len].rstrip()
        + "\n\n[… Phase 2 占位截断；Phase 3 将由 Summarizer 生成摘要]"
    )


def validate_role(role: str) -> str:
    value = (role or "").strip().upper()
    if value not in ROLES:
        raise ValueError(f"role 无效，允许: {', '.join(sorted(ROLES))}")
    return value


def validate_source_type(source_type: str) -> str:
    value = (source_type or "").strip().upper()
    if value not in SOURCE_TYPES:
        raise ValueError(f"source_type 无效，允许: {', '.join(sorted(SOURCE_TYPES))}")
    return value


class DocumentIngestService:
    """将多种输入形态解析为纯文本，UPLOAD 落盘。"""

    def ingest_paste(
        self,
        raw_text: str,
        *,
        title: Optional[str] = None,
    ) -> IngestResult:
        text = sanitize_extracted_text(raw_text or "").strip()
        if not text:
            raise ValueError("粘贴内容不能为空")
        if len(text) > RAW_TEXT_STORE_MAX:
            text = (
                text[:RAW_TEXT_STORE_MAX].rstrip()
                + f"\n\n[… 原文过长，已截断至 {RAW_TEXT_STORE_MAX} 字符入库]"
            )
        return IngestResult(raw_text=text, title=title or "粘贴文本")

    def ingest_bytes(
        self,
        content: bytes,
        *,
        filename: str,
        content_type: Optional[str] = None,
        project_id: Optional[UUID] = None,
    ) -> IngestResult:
        if not content:
            raise ValueError("上传文件为空")
        name = Path(filename or "upload.bin").name
        suffix = Path(name).suffix.lower()
        storage_path = self._save_upload(content, suffix=suffix or ".bin", project_id=project_id)
        text = self.extract_text_from_file(Path(storage_path), filename=name, content_type=content_type)
        text = sanitize_extracted_text(text)
        if not text.strip():
            raise ValueError(
                f"未能从文件「{name}」提取到文本。"
                "若是扫描版 PDF / 图片型 Word，请先转成可选中文字，或改用粘贴。"
            )
        if len(text) > RAW_TEXT_STORE_MAX:
            text = (
                text[:RAW_TEXT_STORE_MAX].rstrip()
                + f"\n\n[… 原文过长，已截断至 {RAW_TEXT_STORE_MAX} 字符入库；完整文件见 storage_path]"
            )
        return IngestResult(
            raw_text=text.strip(),
            storage_path=storage_path,
            original_filename=name,
            content_type=content_type or self._guess_content_type(suffix),
            title=name,
        )

    def extract_text_from_file(
        self,
        path: Path,
        *,
        filename: Optional[str] = None,
        content_type: Optional[str] = None,
    ) -> str:
        """按扩展名/MIME 提取纯文本。"""
        name = (filename or path.name).lower()
        ctype = (content_type or "").lower()
        suffix = path.suffix.lower()

        if suffix in {".md", ".markdown", ".txt"} or "markdown" in ctype or ctype.startswith("text/"):
            return path.read_text(encoding="utf-8", errors="replace")

        if suffix in {".docx", ".dotx", ".docm", ".dotm"} or "wordprocessingml" in ctype:
            from app.services.word_package import (
                friendly_docx_open_error,
                normalize_word_package_inplace,
            )

            normalize_word_package_inplace(path)
            try:
                return pandoc_service.docx_to_markdown(path)
            except ValueError as exc:
                raise ValueError(
                    friendly_docx_open_error(exc, filename or path.name)
                ) from exc

        if suffix == ".pdf" or ctype == "application/pdf":
            return self._extract_pdf(path)

        if suffix == ".pptx" or "presentationml.presentation" in ctype:
            return self._extract_pptx(path)

        if suffix == ".ppt":
            raise ValueError("暂不支持旧版 .ppt，请另存为 .pptx 后上传")

        # 兜底尝试 utf-8
        try:
            return path.read_text(encoding="utf-8")
        except UnicodeDecodeError as exc:
            raise ValueError(f"不支持的文件类型: {name}") from exc

    def _extract_pdf(self, path: Path) -> str:
        try:
            from pypdf import PdfReader
        except ImportError as exc:
            raise ValueError("PDF 解析需要安装 pypdf") from exc

        reader = PdfReader(str(path))
        parts: list[str] = []
        for page in reader.pages:
            text = page.extract_text() or ""
            if text.strip():
                parts.append(text.strip())
        return sanitize_extracted_text("\n\n".join(parts))

    def _extract_pptx(self, path: Path) -> str:
        """提取 PowerPoint (.pptx) 幻灯片文本（含备注）。"""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError as exc:
            raise ValueError("PPTX 解析需要安装 python-pptx") from exc

        prs = Presentation(str(path))
        parts: list[str] = []

        def _shape_text(shape) -> list[str]:  # noqa: ANN001
            chunks: list[str] = []
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                t = (shape.text_frame.text or "").strip()
                if t:
                    chunks.append(t)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    chunks.extend(_shape_text(child))
            if getattr(shape, "has_table", False) and shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [
                        (c.text_frame.text or "").strip().replace("|", "\\|")
                        for c in row.cells
                    ]
                    if any(cells):
                        rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    chunks.append("\n".join(rows))
            return chunks

        for idx, slide in enumerate(prs.slides, start=1):
            slide_bits: list[str] = [f"## Slide {idx}"]
            for shape in slide.shapes:
                slide_bits.extend(_shape_text(shape))
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    slide_bits.append(f"[Notes]\n{notes}")
            body = "\n\n".join(b for b in slide_bits[1:] if b.strip())
            if body.strip():
                parts.append(f"## Slide {idx}\n\n{body.strip()}")
            elif len(slide_bits) == 1:
                continue
            else:
                parts.append(slide_bits[0])

        return sanitize_extracted_text("\n\n".join(parts))

    def _save_upload(
        self,
        content: bytes,
        *,
        suffix: str,
        project_id: Optional[UUID],
    ) -> str:
        root = Path(settings.upload_dir)
        if project_id:
            root = root / str(project_id)
        root.mkdir(parents=True, exist_ok=True)
        # 规范化后缀
        clean_suffix = suffix if suffix.startswith(".") else f".{suffix}"
        path = root / f"{uuid4().hex}{clean_suffix}"
        path.write_bytes(content)
        return str(path)

    @staticmethod
    def _guess_content_type(suffix: str) -> str:
        mapping = {
            ".md": "text/markdown",
            ".markdown": "text/markdown",
            ".txt": "text/plain",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".dotx": "application/vnd.openxmlformats-officedocument.wordprocessingml.template",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".pdf": "application/pdf",
        }
        return mapping.get(suffix.lower(), "application/octet-stream")


document_ingest_service = DocumentIngestService()
