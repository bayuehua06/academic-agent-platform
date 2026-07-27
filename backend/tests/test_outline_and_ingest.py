"""outline_parser / document_ingest 单元测试。"""

from pathlib import Path

from docx import Document

from app.services.document_ingest import (
    document_ingest_service,
    placeholder_summary,
    sanitize_extracted_text,
)
from app.services.outline_parser import parse_docx_outline, parse_markdown_outline
from app.services.pandoc_service import pandoc_service


def test_sanitize_strips_null_bytes():
    dirty = "Hello\x00World\x01\nNext"
    clean = sanitize_extracted_text(dirty)
    assert "\x00" not in clean
    assert "HelloWorld" in clean
    assert "Next" in clean


def test_parse_markdown_outline_levels_and_body():
    md = """# Intro

Alpha paragraph.

## Sub

Beta.

# End
"""
    items = parse_markdown_outline(md)
    assert [i["heading"] for i in items] == ["Intro", "Sub", "End"]
    assert items[0]["level"] == 1
    assert items[1]["level"] == 2
    assert "Alpha" in items[0]["key_points"]
    assert "Beta" in items[1]["key_points"]


def test_parse_docx_outline(tmp_path: Path):
    path = tmp_path / "o.docx"
    doc = Document()
    doc.add_heading("Chapter One", level=1)
    doc.add_paragraph("Details A")
    doc.add_heading("Section Two", level=2)
    doc.add_paragraph("Details B")
    doc.save(str(path))

    items = parse_docx_outline(str(path))
    assert items[0]["heading"] == "Chapter One"
    assert items[0]["level"] == 1
    assert "Details A" in items[0]["key_points"]
    assert items[1]["heading"] == "Section Two"


def test_parse_docx_outline_includes_tables(tmp_path: Path):
    """标题下的表格文字必须进入 key_points（文献向导依赖此字段）。"""
    path = tmp_path / "outline_table.docx"
    doc = Document()
    doc.add_heading("Literature Review", level=1)
    doc.add_paragraph("Intro sentence.")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Theme"
    table.rows[0].cells[1].text = "Focus"
    table.rows[1].cells[0].text = "Platformization"
    table.rows[1].cells[1].text = "Gig economy"
    doc.add_heading("Conclusion", level=1)
    doc.add_paragraph("Wrap up.")
    doc.save(str(path))

    items = parse_docx_outline(str(path))
    by_h = {i["heading"]: i for i in items}
    assert "Literature Review" in by_h
    kp = by_h["Literature Review"]["key_points"]
    assert "Intro sentence" in kp
    assert "Platformization" in kp
    assert "Gig economy" in kp
    assert "Theme" in kp
    assert "Wrap up" in by_h["Conclusion"]["key_points"]
    assert "Platformization" not in by_h["Conclusion"]["key_points"]


def test_placeholder_summary_truncates():
    long = "x" * 5000
    out = placeholder_summary(long, max_len=100)
    assert len(out) < 5000
    assert "Phase 2" in out


def test_ingest_paste_and_markdown_bytes(tmp_path: Path):
    md_path = tmp_path / "a.md"
    md_path.write_text("# Hi\n\nBody", encoding="utf-8")
    text = document_ingest_service.extract_text_from_file(md_path, filename="a.md")
    assert "Hi" in text

    pasted = document_ingest_service.ingest_paste("hello world", title="t")
    assert pasted.raw_text == "hello world"
    assert pasted.title == "t"


def test_ingest_word_template_dotx_content_type(tmp_path: Path):
    """扩展名为 .docx 但 ContentType 是 template.main 时也应能入库。"""
    import zipfile
    from io import BytesIO

    # 造一个最小「模板类型」docx：先写正常文档，再改 Content_Types
    normal = tmp_path / "normal.docx"
    doc = Document()
    doc.add_heading("Chapter From Template", level=1)
    doc.add_paragraph("Body under heading.")
    doc.save(str(normal))

    tmpl = tmp_path / "as_template.docx"
    buf = BytesIO()
    with zipfile.ZipFile(normal, "r") as zin, zipfile.ZipFile(buf, "w") as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            if info.filename == "[Content_Types].xml":
                text = data.decode("utf-8").replace(
                    "wordprocessingml.document.main+xml",
                    "wordprocessingml.template.main+xml",
                )
                data = text.encode("utf-8")
            zout.writestr(info.filename, data)
    tmpl.write_bytes(buf.getvalue())

    # 未规范化前 Document 应失败
    try:
        Document(str(tmpl))
        raised = False
    except ValueError:
        raised = True
    assert raised

    result = document_ingest_service.ingest_bytes(
        tmpl.read_bytes(), filename="outline.docx", project_id=None
    )
    assert "Chapter From Template" in result.raw_text
    items = parse_docx_outline(result.storage_path)
    assert items[0]["heading"] == "Chapter From Template"


def test_docx_table_only_extracts_text(tmp_path: Path):
    """评分表等只有表格、无正文段落的 Word 也应能提取。"""
    path = tmp_path / "rubric.docx"
    doc = Document()
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Criteria"
    table.rows[0].cells[1].text = "Points"
    table.rows[1].cells[0].text = "Challenge Identification"
    table.rows[1].cells[1].text = "/10 pts"
    doc.save(str(path))

    text = pandoc_service.docx_to_markdown(path)
    assert "Challenge Identification" in text
    assert "Criteria" in text


def test_ingest_table_docx_succeeds(tmp_path: Path):
    path = tmp_path / "rubric.docx"
    doc = Document()
    table = doc.add_table(rows=1, cols=1)
    table.rows[0].cells[0].text = "Background of the industry"
    doc.save(str(path))
    result = document_ingest_service.ingest_bytes(
        path.read_bytes(), filename="rubric.docx", project_id=None
    )
    assert "Background of the industry" in result.raw_text


def test_ingest_pptx_extracts_slides(tmp_path: Path):
    from pptx import Presentation

    path = tmp_path / "notes.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Platformization"
    slide.placeholders[1].text = "Key claim about delivery apps."
    slide.notes_slide.notes_text_frame.text = "Remember AUT context."
    prs.save(str(path))
    result = document_ingest_service.ingest_bytes(
        path.read_bytes(), filename="notes.pptx", project_id=None
    )
    assert "Platformization" in result.raw_text
    assert "delivery apps" in result.raw_text
    assert "AUT context" in result.raw_text
